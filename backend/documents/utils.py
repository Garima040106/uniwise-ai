import os
import re
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation


def extract_text_from_file(file_path):
    """
    Extract text from uploaded documents
    Supports: PDF, DOCX, TXT
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == ".pdf":
            return extract_from_pdf(file_path)
        elif ext == ".docx":
            return extract_from_docx(file_path)
        elif ext == ".txt":
            return extract_from_txt(file_path)
        elif ext == ".pptx":
            return extract_from_pptx(file_path)
        else:
            return ""
    except Exception as e:
        print(f"Error extracting text: {e}")
        return ""


def extract_from_pdf(file_path):
    text = ""
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text() or ""
            if page_text:
                text += page_text + "\n"
    return text.strip()


def extract_from_docx(file_path):
    doc = DocxDocument(file_path)
    return "\n".join([para.text for para in doc.paragraphs]).strip()


def extract_from_txt(file_path):
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read().strip()
        except UnicodeDecodeError:
            continue
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()


def extract_from_pptx(file_path):
    presentation = Presentation(file_path)
    lines = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            text = text.strip()
            if text:
                lines.append(text)
    return "\n".join(lines).strip()


def _split_long_text(segment, max_size):
    """
    Split oversized paragraph-like segments into sentence-aware pieces.
    Falls back to whitespace-aware hard splits for very long lines.
    """
    segment = (segment or "").strip()
    if not segment:
        return []
    if len(segment) <= max_size:
        return [segment]

    sentences = re.split(r"(?<=[.!?])\s+", segment)
    pieces = []
    current = ""

    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue

        tentative = f"{current} {sentence}".strip() if current else sentence
        if len(tentative) <= max_size:
            current = tentative
            continue

        if current:
            pieces.append(current)
            current = ""

        if len(sentence) <= max_size:
            current = sentence
            continue

        start = 0
        while start < len(sentence):
            end = min(len(sentence), start + max_size)
            if end < len(sentence):
                split_at = sentence.rfind(" ", start, end)
                if split_at > start + 60:
                    end = split_at
            chunk = sentence[start:end].strip()
            if chunk:
                pieces.append(chunk)
            if end <= start:
                end = start + max_size
            start = end

    if current:
        pieces.append(current)

    return pieces


def chunk_text(text, chunk_size=1000, overlap=100):
    """
    Split text into overlapping, sentence-aware chunks for retrieval.
    """
    clean_text = (text or "").replace("\x00", "").strip()
    if not clean_text:
        return []

    paragraphs = [p.strip() for p in re.split(r"\n{2,}", clean_text) if p.strip()]
    if not paragraphs:
        paragraphs = [clean_text]

    segments = []
    for paragraph in paragraphs:
        segments.extend(_split_long_text(paragraph, chunk_size))

    raw_chunks = []
    current = ""
    for segment in segments:
        tentative = f"{current}\n{segment}".strip() if current else segment
        if len(tentative) <= chunk_size:
            current = tentative
            continue

        if current:
            raw_chunks.append(current.strip())
        current = segment

    if current:
        raw_chunks.append(current.strip())

    if not raw_chunks:
        return []

    overlap = max(0, min(overlap, chunk_size // 2))
    if overlap == 0:
        return raw_chunks

    final_chunks = []
    for index, chunk in enumerate(raw_chunks):
        if index == 0:
            final_chunks.append(chunk)
            continue

        previous_tail = raw_chunks[index - 1][-overlap:].strip()
        if previous_tail and not chunk.startswith(previous_tail):
            merged = f"{previous_tail}\n{chunk}".strip()
        else:
            merged = chunk
        final_chunks.append(merged)

    return final_chunks
