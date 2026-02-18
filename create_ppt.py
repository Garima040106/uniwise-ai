from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(108, 99, 255)
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(left, Inches(3.8), width, Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(108, 99, 255)
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(title, left_title, left_content, right_title, right_content):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(108, 99, 255)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.add_paragraph()
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(67, 233, 123)
    
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(4.2), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.add_paragraph()
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 101, 132)
    
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.space_after = Pt(8)
    
    return slide

# Slide 1: Title
add_title_slide("Uniwise AI", "RAG-Powered University Learning Platform")

# Slide 2: Problem Statement
add_content_slide(
    "The Problem",
    [
        "Students use ChatGPT for homework — but it gives generic answers",
        "Generic AI tools hallucinate facts not in your curriculum",
        "No way to verify if answers come from course materials",
        "Promotes academic dishonesty",
        "Students struggle to retain information from passive reading"
    ]
)

# Slide 3: Our Solution
add_content_slide(
    "Uniwise AI Solution",
    [
        "AI that ONLY answers from professor-uploaded documents",
        "RAG (Retrieval Augmented Generation) prevents hallucinations",
        "Every answer includes source citations",
        "Automatic flashcards & quizzes from your content",
        "Progress tracking and spaced repetition learning",
        "University data isolation — 100% privacy"
    ]
)

# Slide 4: What is RAG?
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "What is RAG?"
p = title_frame.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(108, 99, 255)

# RAG explanation
rag_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
rag_frame = rag_box.text_frame
rag_frame.word_wrap = True

steps = [
    "1. Professor uploads course documents (PDF, DOCX, TXT)",
    "2. System extracts text and creates searchable 'embeddings'",
    "3. Student asks: 'What is photosynthesis?'",
    "4. System searches YOUR documents for relevant paragraphs",
    "5. AI receives ONLY those paragraphs as context",
    "6. AI generates answer using ONLY provided context",
    "7. Answer returned with source citations"
]

for step in steps:
    p = rag_frame.add_paragraph()
    p.text = step
    p.font.size = Pt(16)
    p.space_after = Pt(14)

# Slide 5: Key Features
add_content_slide(
    "Core Features",
    [
        "🤖 Ask AI — Q&A from course materials only",
        "🃏 Auto-Generated Flashcards with spaced repetition",
        "📝 Auto-Generated Quizzes with instant grading",
        "📊 Progress Tracking & Skill Development visualization",
        "👨‍🏫 Professor & Student Roles",
        "🔒 University Data Isolation (multi-tenant)"
    ]
)

# Slide 6: How RAG Prevents Hallucination
add_two_column_slide(
    "RAG vs Generic AI",
    "❌ Generic AI (ChatGPT)",
    [
        "Uses general knowledge from training",
        "May provide outdated information",
        "No source verification",
        "Can make up facts",
        "Not curriculum-specific"
    ],
    "✅ Uniwise AI (RAG)",
    [
        "Only uses YOUR uploaded documents",
        "Always up-to-date with your content",
        "Every answer cites source document",
        "Cannot make up facts",
        "100% curriculum-aligned"
    ]
)

# Slide 7: Architecture Diagram
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "System Architecture"
p = title_frame.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(108, 99, 255)

arch_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(4.5))
arch_frame = arch_box.text_frame
arch_frame.word_wrap = True

arch_text = """React Frontend (Student Interface)
        ↓
Django Backend (Python)
        ↓
┌───────────┬───────────┬───────────┐
│           │           │           │
PostgreSQL  ChromaDB   Ollama AI
(Database)  (Vectors)  (llama3.2)

- PostgreSQL: User data, quizzes, progress
- ChromaDB: Document embeddings (RAG)
- Ollama: Local AI (no cloud, no costs)"""

p = arch_frame.add_paragraph()
p.text = arch_text
p.font.size = Pt(16)
p.font.name = 'Courier New'

# Slide 8: Security & Privacy
add_content_slide(
    "Security & Privacy",
    [
        "🏛️ University Data Isolation — each university has separate database",
        "🔐 Role-Based Access — Professor vs Student permissions",
        "💻 Runs on YOUR server — no data sent to OpenAI/cloud",
        "📜 GDPR/FERPA compliant architecture",
        "✅ All answers cite source documents for verification"
    ]
)

# Slide 9: Technology Stack
add_two_column_slide(
    "Technology Stack",
    "Backend",
    [
        "Django 5.0 (Python framework)",
        "PostgreSQL database",
        "ChromaDB vector store",
        "Ollama (llama3.2:3b — CPU optimized)",
        "Docker containerization"
    ],
    "Frontend",
    [
        "React.js interface",
        "Responsive design",
        "Real-time progress charts",
        "Dark theme UI",
        "Works on any browser"
    ]
)

# Slide 10: Cost Analysis
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Cost Comparison"
p = title_frame.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(108, 99, 255)

cost_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
cost_frame = cost_box.text_frame

costs = [
    ("ChatGPT Plus (500 students)", "$10,000/month", RGBColor(255, 101, 132)),
    ("Chegg (500 students)", "$10,000/month", RGBColor(255, 101, 132)),
    ("Uniwise AI (500 students)", "$100/month", RGBColor(67, 233, 123)),
    ("", "", RGBColor(255, 255, 255)),
    ("Annual Savings", "$118,800/year", RGBColor(67, 233, 123))
]

for item, cost, color in costs:
    p = cost_frame.add_paragraph()
    p.text = f"{item:.<40} {cost}"
    p.font.size = Pt(20)
    p.font.color.rgb = color
    p.font.bold = True if "Savings" in item else False
    p.space_after = Pt(12)

# Slide 11: Use Cases
add_content_slide(
    "Use Cases",
    [
        "👨‍🎓 Students: 24/7 AI tutor using course materials",
        "👨‍🏫 Professors: Auto-generate study materials, reduce repetitive questions",
        "🏛️ Universities: Improve retention, demonstrate innovation",
        "📚 Libraries: Make textbooks searchable with AI Q&A",
        "🎓 Online Courses: Add AI tutoring to MOOCs"
    ]
)

# Slide 12: Benefits Over Competition
add_two_column_slide(
    "Competitive Advantages",
    "Technical",
    [
        "RAG prevents hallucinations",
        "Source citations for every answer",
        "Runs offline (no internet needed)",
        "CPU-optimized (no GPU required)",
        "Open-source technology stack"
    ],
    "Business",
    [
        "One-time setup vs subscription",
        "100x cheaper than ChatGPT",
        "Data stays on university servers",
        "Curriculum-aligned content",
        "Academic integrity built-in"
    ]
)

# Slide 13: POC Status
add_content_slide(
    "Proof of Concept Status",
    [
        "✅ Full RAG implementation complete",
        "✅ Document upload and processing working",
        "✅ AI Q&A, flashcards, quizzes functional",
        "✅ Progress tracking and analytics",
        "✅ Docker deployment ready",
        "✅ Web interface live",
        "🎯 Ready for pilot with 50-100 students"
    ]
)

# Slide 14: Demo
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Live Demo Available"
p = title_frame.paragraphs[0]
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = RGBColor(108, 99, 255)
p.alignment = PP_ALIGN.CENTER

demo_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2))
demo_frame = demo_box.text_frame
demo_frame.text = "Access: http://localhost:3000\n\nGitHub: github.com/[username]/uniwise-ai"
p = demo_frame.paragraphs[0]
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER

# Slide 15: Next Steps
add_content_slide(
    "Next Steps",
    [
        "1. Pilot program with 1-2 courses (50-100 students)",
        "2. Gather feedback and analytics",
        "3. Expand to more courses based on success",
        "4. University-wide rollout",
        "5. Integration with existing LMS (Canvas, Blackboard)",
        "6. Scale to multiple universities"
    ]
)

# Slide 16: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Thank You"
p = title_frame.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = RGBColor(108, 99, 255)
p.alignment = PP_ALIGN.CENTER

contact_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2))
contact_frame = contact_box.text_frame
contact_frame.text = "Questions?\n\nEmail: [your-email]\nGitHub: github.com/[username]/uniwise-ai"
p = contact_frame.paragraphs[0]
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER

# Save
prs.save('/home/garima/uniwise-ai/Uniwise_AI_Presentation.pptx')
print("✅ PowerPoint created: Uniwise_AI_Presentation.pptx")
