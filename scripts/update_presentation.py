from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT_PATH = "../Uniwise_AI_Presentation.pptx"


def set_text_placeholder(placeholder, text, size=24, bold=False, color=(44, 54, 70)):
    tf = placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = RGBColor(*color)


def add_title_slide(prs, title, subtitle=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_ph = slide.shapes.title
    subtitle_ph = slide.placeholders[1]
    set_text_placeholder(title_ph, title, size=40, bold=True, color=(28,40,72))
    if subtitle:
        set_text_placeholder(subtitle_ph, subtitle, size=18, bold=False, color=(90,95,110))


def add_bullet_slide(prs, heading, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    set_text_placeholder(title, heading, size=28, bold=True, color=(28,40,72))

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(80, 90, 110)


def add_two_column_slide(prs, heading, left_items, right_items):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    set_text_placeholder(title, heading, size=28, bold=True, color=(28,40,72))

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(4.0))
    left_tf = left_box.text_frame
    left_tf.clear()
    for it in left_items:
        p = left_tf.add_paragraph()
        p.text = f"• {it}"
        p.font.size = Pt(16)
        p.level = 0
        p.font.color.rgb = RGBColor(70, 80, 100)

    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.0), Inches(4.0))
    right_tf = right_box.text_frame
    right_tf.clear()
    for it in right_items:
        p = right_tf.add_paragraph()
        p.text = f"• {it}"
        p.font.size = Pt(16)
        p.level = 0
        p.font.color.rgb = RGBColor(70, 80, 100)


def build_presentation():
    prs = Presentation()

    # Title
    add_title_slide(prs, "Uniwise AI", "AI-powered adaptive learning for universities")

    # Problem / Value
    add_bullet_slide(prs, "Problem & Value", [
        "Students face cognitive overload and inconsistent pacing.",
        "Educators need scalable, privacy-aware AI tools.",
        "Uniwise AI improves retention with cognitive-aware interventions.",
    ])

    # Features
    add_two_column_slide(prs, "Core Features", [
        "Cognitive load monitoring",
        "Adaptive break enforcement",
        "AI-generated flashcards & quizzes",
        "RAG-based curriculum Q&A",
    ], [
        "Multi-tenant admin portal",
        "Spaced repetition scheduler",
        "Chroma vector DB for embeddings",
        "Ollama local LLM integration",
    ])

    # Architecture
    add_bullet_slide(prs, "Architecture (high-level)", [
        "React frontend ↔ Django REST API",
        "AI Engine: RAG + LLM + CognitiveLoadCalculator",
        "PostgreSQL primary DB + Chroma vector store",
        "Dockerized services for local/prod parity",
    ])

    # Demo flow
    add_bullet_slide(prs, "User Flow / Demo", [
        "Upload document → generate flashcards",
        "Start study session → monitor cognitive load",
        "Receive recommendation → break / review / deep study",
        "Review analytics & progress",
    ])

    # Roadmap (short)
    add_bullet_slide(prs, "Roadmap (short)", [
        "Stabilize tests & CI",
        "Real-time collaboration (WebSockets)",
        "Mobile review app (React Native)",
    ])

    # Contact
    add_bullet_slide(prs, "Contact & Next Steps", [
        "Repo: https://github.com/Garima040106/uniwise-ai",
        "Email: contact@uniwise.ai (configure before public release)",
        "Open issues for features & contributions",
    ])

    # Save
    out = OUTPUT_PATH
    prs.save(out)
    print(f"Presentation written to {out}")


if __name__ == '__main__':
    build_presentation()
